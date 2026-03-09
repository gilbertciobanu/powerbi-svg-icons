#!/usr/bin/env python3
"""
Workflow D — Phase 1: Validate SVG(s) and create PPTX for PowerPoint export.

Usage:
    # Single icon:
    python3 factory/phase1_create_pptx.py [--size WxH] <source_svg> <output_icon_name>

    # Batch mode:
    python3 factory/phase1_create_pptx.py [--size WxH] --batch src1.svg:name1 src2.svg:name2 ...

Arguments:
    source_svg        Filename in icons/ (e.g. ic_fluent_map_24_regular.svg)
                      or absolute path to any SVG
    output_icon_name  Target name in icons/ (e.g. icon-map-regular)
                      without .svg extension

    --batch           Enable batch mode. Each subsequent argument is
                      source_svg:output_icon_name (colon-separated).
                      Creates one slide per icon in a single working.pptx.

    --size WxH        Target resolution in pixels (default: 96).
                      Square shorthand: --size 96 means 96x96.
                      Explicit: --size 96x64.
                      PPTX shape is sized at px/72 inches (1" = 72pt).
                      If aspect ratio differs from source viewBox, exits
                      with an error so the agent can ask the user.

What this script does:
    1. Validates the source SVG(s) for security issues
    2. Embeds them natively in a PPTX at the target resolution
       (default 96px = 1.333") → factory/working.pptx
       (one slide per icon in batch mode)

Next step (Windows):
    Single: Run ExportIconAsSVG macro in PowerPoint (exports output.svg)
    Batch:  Run ExportBatchSVGs macro in PowerPoint (exports output_1.svg, output_2.svg, ...)

    Then tell the agent to run phase 2.
"""

import sys, os, re, json, zipfile, subprocess, datetime
from pathlib import Path

FACTORY_DIR = Path(__file__).parent
PROJECT_DIR = FACTORY_DIR.parent
ICONS_DIR   = PROJECT_DIR / "icons"
STATE_FILE  = FACTORY_DIR / ".workflow_state"


# ── Security validation ──────────────────────────────────────────────────────

def validate_svg(path: Path) -> list[str]:
    """Return list of security warnings. Empty = clean."""
    content = path.read_text(encoding="utf-8", errors="replace")
    issues = []

    if re.search(r'<script', content, re.IGNORECASE):
        issues.append("<script> tag detected — could execute code")

    for attr in ("href", "src", "xlink:href"):
        for url in re.findall(rf'{attr}=["\']([^"\']+)["\']', content):
            if url.startswith("http") and "w3.org" not in url:
                issues.append(f"External URL in {attr}: {url}")

    if re.search(r'<image[^>]+href=["\']https?://', content):
        issues.append("External <image> href found")

    if re.search(r'<use[^>]+href=["\']https?://', content):
        issues.append("External <use> reference (will not resolve in sandbox)")

    return issues


# ── Parse --size flag ────────────────────────────────────────────────────────

DEFAULT_SIZE_PX = 96  # 96px = 96/72 = 1.333"

def parse_size(size_str: str) -> tuple[int, int]:
    """Parse size string like '96' (square) or '96x64' into (w, h) pixels."""
    if "x" in size_str.lower():
        parts = size_str.lower().split("x", 1)
        return int(parts[0]), int(parts[1])
    val = int(size_str)
    return val, val


def get_viewbox(svg_path: Path) -> tuple[float, float]:
    """Extract (width, height) from an SVG's viewBox attribute."""
    content = svg_path.read_text(encoding="utf-8", errors="replace")
    m = re.search(r'viewBox=["\']([^"\']+)["\']', content, re.IGNORECASE)
    if m:
        parts = m.group(1).split()
        if len(parts) == 4:
            return float(parts[2]), float(parts[3])
    # Fallback: try width/height attributes
    wm = re.search(r'\bwidth=["\'](\d+)', content)
    hm = re.search(r'\bheight=["\'](\d+)', content)
    if wm and hm:
        return float(wm.group(1)), float(hm.group(1))
    return 0, 0


def check_proportionality(src_path: Path, target_w: int, target_h: int) -> str | None:
    """Return an error message if target aspect ratio doesn't match source, else None."""
    vb_w, vb_h = get_viewbox(src_path)
    if vb_w == 0 or vb_h == 0:
        return None  # Can't check, allow
    src_ratio = vb_w / vb_h
    tgt_ratio = target_w / target_h
    if abs(src_ratio - tgt_ratio) / max(src_ratio, tgt_ratio) > 0.01:
        return (
            f"{src_path.name}: source aspect ratio {vb_w:.0f}:{vb_h:.0f} "
            f"({src_ratio:.3f}) ≠ target {target_w}:{target_h} ({tgt_ratio:.3f})"
        )
    return None


# ── PPTX creation (single slide) ────────────────────────────────────────────

def create_pptx(svg_path: Path, output_pptx: Path, size_px: tuple[int, int] = None):
    """Embed a single SVG natively in PPTX on a blank slide."""
    _create_pptx_multi([svg_path], output_pptx, size_px=size_px)


# ── PPTX creation (multi-slide) ─────────────────────────────────────────────

def _create_pptx_multi(svg_paths: list[Path], output_pptx: Path,
                        size_px: tuple[int, int] = None):
    """Embed one or more SVGs natively in PPTX (one slide per SVG)."""
    from pptx import Presentation
    from pptx.util import Inches
    from lxml import etree

    if size_px is None:
        size_px = (DEFAULT_SIZE_PX, DEFAULT_SIZE_PX)
    w_inches = size_px[0] / 72
    h_inches = size_px[1] / 72

    n = len(svg_paths)
    png_tmps = []

    # Create slides with PNG fallbacks
    prs = Presentation()
    for i, svg_path in enumerate(svg_paths):
        png_tmp = FACTORY_DIR / f"_fallback_{i+1}.png"
        subprocess.run(
            ["convert", "-background", "none", "-density", "384",
             "-resize", "384x384", str(svg_path), str(png_tmp)],
            check=True
        )
        png_tmps.append(png_tmp)

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.add_picture(str(png_tmp), Inches(1), Inches(1),
                                 Inches(w_inches), Inches(h_inches))

    raw_pptx = FACTORY_DIR / "_raw.pptx"
    prs.save(str(raw_pptx))

    # Read all SVG bytes
    svg_bytes_list = [p.read_bytes() for p in svg_paths]

    A_NS     = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ASVG_NS  = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
    R_NS     = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

    with zipfile.ZipFile(str(raw_pptx), "r") as zin, \
         zipfile.ZipFile(str(output_pptx), "w", zipfile.ZIP_DEFLATED) as zout:

        for item in zin.infolist():
            data = zin.read(item.filename)

            # Add SVG content type (once)
            if item.filename == "[Content_Types].xml":
                ct = data.decode("utf-8")
                if "svg+xml" not in ct:
                    ct = ct.replace(
                        "</Types>",
                        '<Default Extension="svg" ContentType="image/svg+xml"/></Types>'
                    )
                data = ct.encode("utf-8")

            else:
                # Check if this is a slide .rels or slide .xml for any slide
                for slide_idx in range(1, n + 1):
                    svg_rid = f"rId{98 + slide_idx}"  # rId99, rId100, etc.
                    svg_media = f"image{slide_idx}.svg"

                    if item.filename == f"ppt/slides/_rels/slide{slide_idx}.xml.rels":
                        rels = data.decode("utf-8")
                        new_rel = (
                            f'<Relationship Id="{svg_rid}" Type="{REL_TYPE}"'
                            f' Target="../media/{svg_media}"/>'
                        )
                        rels = rels.replace("</Relationships>", new_rel + "</Relationships>")
                        data = rels.encode("utf-8")
                        break

                    elif item.filename == f"ppt/slides/slide{slide_idx}.xml":
                        root = etree.fromstring(data)
                        blip = root.find(f".//{{{A_NS}}}blip")
                        if blip is not None:
                            extLst  = etree.SubElement(blip,   f"{{{A_NS}}}extLst")
                            ext     = etree.SubElement(extLst,  f"{{{A_NS}}}ext")
                            ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
                            svgBlip = etree.SubElement(ext, f"{{{ASVG_NS}}}svgBlip")
                            svgBlip.set(f"{{{R_NS}}}embed", svg_rid)
                        data = etree.tostring(root, xml_declaration=True,
                                              encoding="UTF-8", standalone=True)
                        break

            zout.writestr(item, data)

        # Write all SVG media files
        for i, svg_bytes in enumerate(svg_bytes_list, start=1):
            zout.writestr(f"ppt/media/image{i}.svg", svg_bytes)

    raw_pptx.unlink(missing_ok=True)
    for p in png_tmps:
        p.unlink(missing_ok=True)


# ── Resolve source path ──────────────────────────────────────────────────────

def resolve_source(src_arg: str) -> Path:
    """Resolve a source SVG argument to an absolute Path."""
    src = Path(src_arg)
    if not src.is_absolute():
        src = ICONS_DIR / src_arg
    return src


# ── Argument parsing ─────────────────────────────────────────────────────────

def parse_args(argv: list[str]) -> dict:
    """Parse CLI arguments into a structured dict.

    Returns:
        {
            'mode': 'single' | 'batch',
            'size': (w_px, h_px),
            'pairs': [(Path, icon_name), ...],   # batch
            'source': Path,                       # single
            'icon_name': str,                     # single
        }
    """
    args = argv[1:]  # skip script name
    size_px = (DEFAULT_SIZE_PX, DEFAULT_SIZE_PX)
    batch = False
    positional = []

    i = 0
    while i < len(args):
        if args[i] == "--size" and i + 1 < len(args):
            size_px = parse_size(args[i + 1])
            i += 2
        elif args[i] == "--batch":
            batch = True
            i += 1
        else:
            positional.append(args[i])
            i += 1

    result = {"size": size_px}

    if batch:
        result["mode"] = "batch"
        pairs = []
        for arg in positional:
            if ":" not in arg:
                print(f"ERROR: Invalid batch argument '{arg}'. Expected source.svg:icon-name")
                sys.exit(1)
            src_arg, icon_name = arg.split(":", 1)
            src = resolve_source(src_arg)
            if not src.exists():
                print(f"ERROR: Source SVG not found: {src}")
                sys.exit(1)
            pairs.append((src, icon_name))
        if not pairs:
            print("ERROR: --batch requires at least one source:icon_name pair.")
            sys.exit(1)
        result["pairs"] = pairs
    else:
        result["mode"] = "single"
        if len(positional) < 2:
            print(__doc__)
            sys.exit(1)
        src = resolve_source(positional[0])
        if not src.exists():
            print(f"ERROR: Source SVG not found: {src}")
            sys.exit(1)
        result["source"] = src
        result["icon_name"] = positional[1]

    return result


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    cfg = parse_args(sys.argv)
    size_px = cfg["size"]
    size_label = f"{size_px[0]}x{size_px[1]}"
    w_in = size_px[0] / 72
    h_in = size_px[1] / 72

    # ── Batch mode ───────────────────────────────────────────────────────────
    if cfg["mode"] == "batch":
        pairs = cfg["pairs"]
        n = len(pairs)

        print(f"\n{'='*60}")
        print(f"  Batch mode — {n} icon(s)  [target {size_label} px = {w_in:.3f}\"×{h_in:.3f}\"]")
        print(f"{'='*60}\n")

        # Print summary table
        print(f"  {'Slide':<6} {'Source':<45} {'Output'}")
        print(f"  {'─'*6} {'─'*45} {'─'*30}")
        for i, (src, name) in enumerate(pairs, start=1):
            print(f"  {i:<6} {src.name:<45} {name}.svg")

        # Step 1: Validate all SVGs + proportionality check
        print(f"\n[1/2] Validating {n} SVGs...")
        has_errors = False
        for src, name in pairs:
            issues = validate_svg(src)
            if issues:
                has_errors = True
                print(f"  WARNING: {src.name}")
                for w in issues:
                    print(f"         - {w}")
            else:
                print(f"  ✓ {src.name}")

            prop_err = check_proportionality(src, size_px[0], size_px[1])
            if prop_err:
                print(f"  ERROR: {prop_err}")
                has_errors = True

        if has_errors:
            # Check if any proportionality errors (hard stop)
            prop_errors = [
                check_proportionality(src, size_px[0], size_px[1])
                for src, _ in pairs
            ]
            if any(e is not None for e in prop_errors):
                print("\n  Aspect ratio mismatch detected. Aborting.")
                print("  Ask the user to choose a proportional target size.")
                sys.exit(1)
            print("\n  Some SVGs have security warnings (see above). Continuing...")

        # Step 2: Create multi-slide working.pptx
        pptx_out = FACTORY_DIR / "working.pptx"
        print(f"\n[2/2] Creating PPTX with {n} slides ({size_label} px per shape)...")
        _create_pptx_multi([src for src, _ in pairs], pptx_out, size_px=size_px)
        print(f"      OK — factory/working.pptx ({n} slides)")

        # Save batch state as JSON
        state = {
            "mode": "batch",
            "target_size": size_label,
            "icons": [
                {"slide": i + 1, "source": src.name, "icon_name": name}
                for i, (src, name) in enumerate(pairs)
            ],
            "created": datetime.datetime.now().isoformat(timespec="seconds"),
        }
        STATE_FILE.write_text(json.dumps(state, indent=2), encoding="utf-8")

        print(f"""
{'='*60}
  Phase 1 complete — {n} icons queued.
  Shape size: {size_label} px ({w_in:.3f}\"×{h_in:.3f}\")

  Next step (Windows):
    Open macro-seed.pptm → Alt+F8 → run ExportBatchSVGs

    The macro will export output_1.svg … output_{n}.svg
    to the factory/ folder.

  Then tell the agent to run phase 2.
{'='*60}
""")

    # ── Single-icon mode ─────────────────────────────────────────────────────
    else:
        src = cfg["source"]
        icon_name = cfg["icon_name"]

        # Proportionality check
        prop_err = check_proportionality(src, size_px[0], size_px[1])
        if prop_err:
            print(f"ERROR: {prop_err}")
            print("Ask the user to choose a proportional target size.")
            sys.exit(1)

        # Check output name doesn't already exist
        out_svg = ICONS_DIR / f"{icon_name}.svg"
        if out_svg.exists():
            print(f"WARNING: {out_svg.name} already exists in icons/")
            confirm = input("Overwrite when phase 2 runs? (yes/no): ").strip().lower()
            if confirm != "yes":
                print("Aborted.")
                sys.exit(0)

        print(f"\nSource : {src.name}")
        print(f"Output : {icon_name}.svg")
        print(f"Size   : {size_label} px ({w_in:.3f}\"×{h_in:.3f}\")")

        # Step 1: Validate
        print("\n[1/2] Validating SVG...")
        issues = validate_svg(src)
        if issues:
            print("WARNING: Security issues found:")
            for w in issues:
                print(f"     - {w}")
            confirm = input("Proceed anyway? (yes/no): ").strip().lower()
            if confirm != "yes":
                print("Aborted.")
                sys.exit(1)
        else:
            print("      OK - Clean")

        # Step 2: Create working.pptx
        pptx_out = FACTORY_DIR / "working.pptx"
        print(f"\n[2/2] Creating PPTX with SVG at {size_label} px...")
        create_pptx(src, pptx_out, size_px=size_px)
        print(f"      OK - factory/working.pptx")

        # Save state for phase 2 (flat format — backward compatible)
        STATE_FILE.write_text(
            f"icon_name={icon_name}\n"
            f"source={src.name}\n"
            f"target_size={size_label}\n"
            f"working_file=working.pptx\n"
            f"created={datetime.datetime.now().isoformat(timespec='seconds')}\n"
        )

        print(f"""
=========================================
  Phase 1 complete.

  Next step (Windows):
    Open macro-seed.pptm → Alt+F8 → run ExportIconAsSVG
    Select the factory/ folder when prompted.

  Then tell the agent to run phase 2.
=========================================
""")


if __name__ == "__main__":
    main()
